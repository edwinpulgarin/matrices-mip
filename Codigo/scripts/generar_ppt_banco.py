# -*- coding: utf-8 -*-
"""Genera el PPT de Banco Caja Social — 4 slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy
from lxml import etree

# ── COLORES ──────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x2B, 0x6E)   # fondo oscuro títulos
BLUE      = RGBColor(0x10, 0x5F, 0xC0)   # azul medio
LIGHTBLUE = RGBColor(0xE8, 0xF4, 0xFD)   # fondo card claro
CYAN      = RGBColor(0x00, 0xB4, 0xE0)   # acento cyan
DARKCARD  = RGBColor(0x10, 0x4E, 0xA8)   # card inferior azul oscuro
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x0D, 0x1B, 0x40)
GRAY      = RGBColor(0x44, 0x44, 0x66)
RED_LABEL = RGBColor(0xCC, 0x22, 0x22)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # blank

# ── HELPERS ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, l, t, w, h, fill=None, radius=Emu(100000),
                     line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(5, l, t, w, h)  # rounded rectangle
    shape.adjustments[0] = 0.05
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def txb(slide, text, l, t, w, h,
        size=Pt(11), bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, wrap=True,
        italic=False, line_spacing=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    if line_spacing:
        from pptx.util import Pt as pt2
        p.line_spacing = line_spacing
    return tb

def rich_para(tf, segments, align=PP_ALIGN.LEFT, space_before=None):
    """segments = list of (text, bold, color, size, italic)"""
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = space_before
    for text, bold, color, size, italic in segments:
        r = p.add_run()
        r.text = text
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.size = size
        r.font.italic = italic
    return p

def set_first_para(tf, segments, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0]
    p.alignment = align
    p.clear()
    for text, bold, color, size, italic in segments:
        r = p.add_run()
        r.text = text
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.size = size
        r.font.italic = italic
    return p

# ── SLIDE 1 ───────────────────────────────────────────────────────────────────
# "La relación de los clientes con sus bancos..."
s1 = prs.slides.add_slide(blank_layout)

# White background
add_rect(s1, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)

# Title
txb(s1,
    "La relación de los clientes con sus bancos está cada vez más\n"
    "fragmentada, mientras que la lealtad se ha erosionado",
    Inches(.45), Inches(.2), Inches(12.4), Inches(1.1),
    size=Pt(22), bold=True, color=NAVY)

# ── 3 CARDS ──
card_tops    = [Inches(1.45)] * 3
card_heights = [Inches(5.6)]  * 3
card_widths  = [Inches(3.9)]  * 3
card_lefts   = [Inches(0.35), Inches(4.55), Inches(8.75)]

body_top    = Inches(1.45)
body_h      = Inches(3.35)
bottom_top  = Inches(4.9)
bottom_h    = Inches(2.15)

card_data = [
    {
        "num": "1",
        "title": "Los clientes tienen un mayor\nnúmero de relaciones bancarias",
        "body": ("Varios bancos y empresas fintech se dirigen a los clientes "
                 "ofreciéndoles una amplia gama de propuestas de valor, ",
                 "lo que fragmenta la relación con el cliente",
                 ""),
        "stat": ("Aproximadamente el ", "70%\nde los clientes",
                 " tienen más de una relación bancaria, cifra que ha aumentado.\n"
                 "23 puntos porcentuales en los últimos 4 años"),
    },
    {
        "num": "2",
        "title": "Los clientes tienen mayores\nexpectativas y están\ndispuestos a cambiar de\nbanco principal",
        "body_rich": [
            ("Los clientes tienen ", False),
            ("mayores expectativas en cuanto", True),
            (" a la ", False),
            ("experiencia del usuario, el servicio al cliente y la cartera de productos,", True),
            (" y están ", False),
            ("dispuestos a cambiarse a su banco principal si no están satisfechos", True),
        ],
        "stat": ("Aproximadamente el ", "40%\nde los clientes",
                 " consideraría cambiar de banco"),
    },
    {
        "num": "3",
        "title": "Las fintech están ganando\nterreno en ciertos productos\nde nicho",
        "body_rich": [
            ("El panorama competitivo está cambiando rápidamente; los bancos tradicionales y las", True),
            (" fintech se centran en la población bancarizada con ofertas innovadoras y modelos operativos digitales", False),
        ],
        "stat": ("Aproximadamente ", "2.300\nempresas fintech",
                 " en Latinoamérica, principalmente en tecnología de pagos, remesas y banca"),
    },
]

for i, (cl, cd) in enumerate(zip(card_lefts, card_data)):
    cw = card_widths[i]

    # Card border (rounded)
    border = add_rounded_rect(s1, cl, body_top, cw, Inches(5.6),
                              fill=WHITE,
                              line_color=RGBColor(0xB0, 0xCC, 0xEE),
                              line_width=Pt(1.2))

    # Number badge
    badge = add_rect(s1, cl + Inches(0.18), body_top + Inches(0.18),
                     Inches(0.38), Inches(0.38), fill=NAVY)
    txb(s1, cd["num"],
        cl + Inches(0.18), body_top + Inches(0.13),
        Inches(0.38), Inches(0.45),
        size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Card title
    txb(s1, cd["title"],
        cl + Inches(0.65), body_top + Inches(0.12),
        cw - Inches(0.82), Inches(0.95),
        size=Pt(11.5), bold=True, color=NAVY)

    # Card body text
    tb = slide_txb = s1.shapes.add_textbox(
        cl + Inches(0.18), body_top + Inches(1.2),
        cw - Inches(0.36), Inches(1.85))
    tf = tb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT

    if "body_rich" in cd:
        for j, (txt, bld) in enumerate(cd["body_rich"]):
            r = p0.add_run()
            r.text = txt
            r.font.size = Pt(10.5)
            r.font.bold = bld
            r.font.color.rgb = GRAY if not bld else BLACK
    else:
        b1, b2, b3 = cd["body"]
        r = p0.add_run(); r.text = b1
        r.font.size = Pt(10.5); r.font.color.rgb = GRAY
        r2 = p0.add_run(); r2.text = b2
        r2.font.size = Pt(10.5); r2.font.bold = True; r2.font.color.rgb = BLACK

    # Bottom dark box
    add_rect(s1, cl, bottom_top, cw, bottom_h, fill=DARKCARD)

    # Stat text
    tb2 = s1.shapes.add_textbox(cl + Inches(0.2), bottom_top + Inches(0.15),
                                 cw - Inches(0.3), bottom_h - Inches(0.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    pre, bold_part, post = cd["stat"]
    r_pre  = p2.add_run(); r_pre.text = pre
    r_pre.font.size  = Pt(10); r_pre.font.color.rgb = WHITE
    r_bld  = p2.add_run(); r_bld.text = bold_part
    r_bld.font.size  = Pt(11); r_bld.font.bold = True; r_bld.font.color.rgb = WHITE
    r_post = p2.add_run(); r_post.text = post
    r_post.font.size = Pt(10); r_post.font.color.rgb = WHITE

# Footer
txb(s1,
    "Fuente: Datos basados en una encuesta sobre banca en mercados emergentes – octubre de 2023. Statista",
    Inches(0.35), Inches(7.1), Inches(8), Inches(0.3),
    size=Pt(7.5), color=GRAY, italic=True)
txb(s1, "Banco Caja Social  |  3",
    Inches(10.5), Inches(7.1), Inches(2.7), Inches(0.3),
    size=Pt(8), color=NAVY, align=PP_ALIGN.RIGHT, bold=True)


# ── SLIDE 2 ───────────────────────────────────────────────────────────────────
# "Tenemos 5 principios rectores..."
s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, SLIDE_W, SLIDE_H, fill=BLUE)

# Title
txb(s2,
    "Tenemos 5 principios rectores que sientan las bases de nuestra\nmetodología de principalidad",
    Inches(.45), Inches(.18), Inches(4.2), Inches(1.35),
    size=Pt(18), bold=True, color=WHITE)

# Center circle
def add_oval(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(9, l, t, w, h)  # oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

CX = Inches(6.0)
CY = Inches(4.0)
R_outer = Inches(1.55)
R_inner = Inches(1.2)

# Outer ring
outer = add_oval(s2, CX - R_outer, CY - R_outer, R_outer*2, R_outer*2,
                 fill=RGBColor(0x08, 0x3A, 0x8A))
inner = add_oval(s2, CX - R_inner, CY - R_inner, R_inner*2, R_inner*2,
                 fill=RGBColor(0x0D, 0x2B, 0x6E))

txb(s2, "Principios\nrectores",
    CX - R_inner + Inches(0.1), CY - Inches(0.4),
    R_inner * 2 - Inches(0.2), Inches(0.9),
    size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 5 satellite circles + labels
import math

principles = [
    ("Establecer principalidad\ncomo eje de gestión del\nvalor del cliente",   math.pi/2),       # top
    ("Crear una visión única,\nintegrada y coherente\nde cada cliente",  math.pi/2 - 2*math.pi/5),   # upper right
    ("Usar datos y analítica\npara entender necesidades\ny personalizar experiencias",   math.pi/2 - 4*math.pi/5),  # lower right
    ("Adoptar un enfoque\ncentrado en cliente,\nno en producto",   math.pi/2 + 4*math.pi/5),   # lower left
    ("Crear una gobernanza\nestructurada para supervisar\nlos resultados y gestionar la\nresponsabilidad en todo el\nbanco",  math.pi/2 + 2*math.pi/5),    # upper left
]

R_orbit = Inches(2.55)
r_sat   = Inches(0.42)

for label, angle in principles:
    sx = CX + R_orbit * math.cos(angle)
    sy = CY - R_orbit * math.sin(angle)

    # satellite circle
    sat = add_oval(s2, sx - r_sat, sy - r_sat, r_sat*2, r_sat*2,
                   fill=RGBColor(0xA8, 0xC8, 0xF0))

    # line from center to satellite (thin)
    from pptx.util import Pt as Pt2
    conn = s2.shapes.add_shape(9, sx - Inches(0.05), sy - Inches(0.05),
                               Inches(0.1), Inches(0.1))
    conn.fill.background(); conn.line.fill.background()

    # label placement
    lw = Inches(2.2)
    lh = Inches(1.1)
    if math.cos(angle) > 0.2:
        lx = sx + r_sat + Inches(0.12)
    elif math.cos(angle) < -0.2:
        lx = sx - r_sat - lw - Inches(0.12)
    else:
        lx = sx - lw/2

    if math.sin(angle) > 0.3:
        ly = sy - lh - Inches(0.3)
    elif math.sin(angle) < -0.3:
        ly = sy + r_sat + Inches(0.1)
    else:
        ly = sy - lh/2

    al = PP_ALIGN.LEFT if math.cos(angle) > 0 else PP_ALIGN.RIGHT if math.cos(angle) < -0.2 else PP_ALIGN.CENTER
    txb(s2, label, lx, ly, lw, lh,
        size=Pt(9.5), color=WHITE, align=al)

txb(s2, "Banco Caja Social  |  6",
    Inches(10.5), Inches(7.1), Inches(2.7), Inches(0.3),
    size=Pt(8), color=WHITE, align=PP_ALIGN.RIGHT, bold=True)


# ── SLIDE 3 ───────────────────────────────────────────────────────────────────
# "Para cada cliente, creamos el indicador de principalidad..."
s3 = prs.slides.add_slide(blank_layout)
add_rect(s3, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)

txb(s3,
    "2. Para cada cliente, creamos el indicador de principalidad\nutilizando una escala única de 0 a 100",
    Inches(.4), Inches(.15), Inches(12.5), Inches(1.0),
    size=Pt(22), bold=True, color=NAVY)

# Labels
txb(s3, "EJEMPLO DE CLIENTE SANITIZADO – ILUSTRATIVO",
    Inches(.4), Inches(1.15), Inches(4.5), Inches(0.3),
    size=Pt(8), bold=True, color=RED_LABEL)

txb(s3, "Escala del principalidad:",
    Inches(4.5), Inches(1.15), Inches(4), Inches(0.28),
    size=Pt(10), bold=True, color=NAVY)

# Gradient bar 0-100
bar_l, bar_t, bar_w, bar_h = Inches(4.5), Inches(1.5), Inches(8.5), Inches(0.22)
# Simulate gradient with multiple rects
steps = 40
for k in range(steps):
    frac = k / steps
    r = int(0x10 + frac * (0x00 - 0x10))
    g = int(0x5F + frac * (0xB4 - 0x5F))
    b = int(0xC0 + frac * (0xE0 - 0xC0))
    add_rect(s3,
             bar_l + int(bar_w * k/steps),
             bar_t,
             int(bar_w / steps) + 1, bar_h,
             fill=RGBColor(r, g, b))
txb(s3, "0",   Inches(4.5),  Inches(1.76), Inches(.3), Inches(.25), size=Pt(9), color=NAVY)
txb(s3, "100", Inches(12.7), Inches(1.76), Inches(.5), Inches(.25), size=Pt(9), color=NAVY)

# Column headers
col_labels = ["principalidad\nbajo o nulo", "Clientes\nactivos", "Clientes\nleales"]
col_centers = [Inches(5.5), Inches(8.2), Inches(11.2)]
for lbl, cx in zip(col_labels, col_centers):
    txb(s3, lbl, cx - Inches(1.2), Inches(2.1), Inches(2.4), Inches(0.5),
        size=Pt(10.5), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Table
rows_data = [
    ("Indicadores mensuales",              None,      "Estrategia",              None,       None),
    ("Gasto mensual (R$)",                 ["50","500","2.000"]),
    ("Uso de la tarjeta (#transacciones)", ["1","5","45"]),
    ("Transferencias (#transacciones)",    ["0","3","17"]),
    ("Inversiones (#operaciones)",         ["0","1","2"]),
    ("Servicio Nacional de Salud",         ["40%","60%","75%"]),
    ("Tasa de abandono anualizada (%)",    ["7%","2%","1%"]),
]

TBL_L = Inches(.4)
TBL_T = Inches(2.7)
TBL_W = Inches(12.6)
ROW_H = Inches(0.42)
COL_W = [Inches(3.5), Inches(2.8), Inches(2.8), Inches(3.2)]

row_labels = [
    "Gasto mensual (R$)",
    "Uso de la tarjeta (#transacciones)",
    "Transferencias (#transacciones)",
    "Inversiones (#operaciones)",
    "Servicio Nacional de Salud",
    "Tasa de abandono anualizada (%)",
]
row_values = [
    ["50","500","2.000"],
    ["1","5","45"],
    ["0","3","17"],
    ["0","1","2"],
    ["40%","60%","75%"],
    ["7%","2%","1%"],
]

# Header row
add_rect(s3, TBL_L, TBL_T, TBL_W, ROW_H, fill=WHITE,
         line_color=RGBColor(0xCC,0xDD,0xEE), line_width=Pt(.5))
txb(s3, "Indicadores mensuales", TBL_L + Inches(.08), TBL_T + Inches(.08),
    COL_W[0], ROW_H - Inches(.08), size=Pt(10), bold=True, color=NAVY)

for r_i, (lbl, vals) in enumerate(zip(row_labels, row_values)):
    ry = TBL_T + ROW_H * (r_i + 1)
    bg = RGBColor(0xF8,0xFB,0xFF) if r_i % 2 == 0 else WHITE
    add_rect(s3, TBL_L, ry, TBL_W, ROW_H, fill=bg,
             line_color=RGBColor(0xCC,0xDD,0xEE), line_width=Pt(.4))
    txb(s3, lbl, TBL_L + Inches(.1), ry + Inches(.08),
        COL_W[0] - Inches(.1), ROW_H, size=Pt(10), color=BLACK)
    for ci, (val, cx) in enumerate(zip(vals, col_centers)):
        pill = add_rounded_rect(s3, cx - Inches(.55), ry + Inches(.07),
                                Inches(1.1), Inches(.28), fill=CYAN)
        txb(s3, val, cx - Inches(.55), ry + Inches(.04),
            Inches(1.1), Inches(.35), size=Pt(10), bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

# Strategy section
STRAT_T = TBL_T + ROW_H * 7 + Inches(.1)
txb(s3, "Estrategia", TBL_L, STRAT_T, Inches(1.5), Inches(.3),
    size=Pt(10), bold=True, color=NAVY)

strat_cols = [
    "• Nuevos clientes: Proceso de integración\n• Clientes en riesgo de abandono: recuperación personalizada\n• Monoproducto: Expansión de la relación a transaccional",
    "• Jerarquía de palancas de incremento del proyecto\n• Definición de acciones comerciales\n• Mejor canal de contacto/oferta",
    "• Comunicación proactiva de beneficios\n• alertas de pérdida de relación y pronóstico del ciclo de vida",
]
scol_lefts = [Inches(.4), Inches(4.8), Inches(8.9)]
for txt, sl in zip(strat_cols, scol_lefts):
    txb(s3, txt, sl, STRAT_T + Inches(.32), Inches(4.1), Inches(.95),
        size=Pt(8.5), color=BLACK)

txb(s3, "Banco Caja Social  |  10",
    Inches(10.5), Inches(7.1), Inches(2.7), Inches(0.3),
    size=Pt(8), color=NAVY, align=PP_ALIGN.RIGHT, bold=True)


# ── SLIDE 4 ───────────────────────────────────────────────────────────────────
# "De este modo, podemos evaluar y comparar diferentes indicadores entre grupos"
s4 = prs.slides.add_slide(blank_layout)
add_rect(s4, 0, 0, SLIDE_W, SLIDE_H, fill=BLUE)

txb(s4,
    "1. De este modo, podemos evaluar y comparar diferentes\nindicadores entre grupos",
    Inches(.4), Inches(.15), Inches(12.5), Inches(1.0),
    size=Pt(22), bold=True, color=WHITE)

txb(s4, "EJEMPLO DE CLIENTE SANITIZADO – ILUSTRATIVO",
    Inches(.4), Inches(1.1), Inches(5), Inches(0.25),
    size=Pt(8), bold=True, color=RGBColor(0xFF,0xCC,0xCC))

# Group labels header
txb(s4, "Entradas para agrupación",
    Inches(3.5), Inches(1.12), Inches(5.5), Inches(0.22),
    size=Pt(9), color=WHITE, bold=True)

# Column group headers
txb(s4, "Datos sociodemográficos",
    Inches(2.0), Inches(1.38), Inches(3.0), Inches(0.22),
    size=Pt(8.5), color=WHITE, bold=True)
txb(s4, "Necesidades financieras y plazos",
    Inches(5.2), Inches(1.38), Inches(3.2), Inches(0.22),
    size=Pt(8.5), color=WHITE, bold=True)
txb(s4, "Características del clúster",
    Inches(8.55), Inches(1.38), Inches(4.3), Inches(0.22),
    size=Pt(8.5), color=WHITE, bold=True)

txb(s4, "Cifras en base 100",
    Inches(.1), Inches(1.62), Inches(1.8), Inches(.22),
    size=Pt(7.5), color=WHITE, italic=True)

# Sub-column headers
sub_cols = [
    ("Promedio edad\nAños",    Inches(2.0)),
    ("Promedio ingreso\nR$ k", Inches(2.72)),
    ("Familia ingreso\nR$ k",  Inches(3.45)),
    ("% de operaciones\nbancarias diarias", Inches(4.45)),
    ("Balance de gastos\n%",   Inches(5.65)),
    ("Construcción del patrimonio\n%", Inches(6.75)),
    ("clientes\ndigitales %",  Inches(8.1)),
    ("Ingresos por cliente\nR$/mes", Inches(9.1)),
    ("Servicio Nacional\nde Salud", Inches(10.2)),
    ("Nuevo competidor\n% con interacción", Inches(11.3)),
]
for lbl, lft in sub_cols:
    txb(s4, lbl, lft, Inches(1.62), Inches(.95), Inches(.52),
        size=Pt(7), color=WHITE, align=PP_ALIGN.CENTER)

# Cluster rows
clusters = [
    ("Grupo 1",   "22","30","75", "98%", "12%","27%", "91%","-3",  "59","27%"),
    ("Clúster 2", "22","23","45", "100%","17%","32%", "71%","-6",  "71","32%"),
    ("Clúster 3", "37","68","80", "57%", "24%","91%", "83%","33",  "75","91%"),
    ("Grupo 4",   "39","33","55", "75%", "33%","73%", "68%","17",  "74","73%"),
    ("Grupo 5",   "41","100","100","44%","17%","91%", "100%","100","75","91%"),
    ("Clúster 6", "43","55","75", "51%", "40%","100%","91%","88",  "71","100%"),
    ("Clúster 7", "46","98","88", "42%", "29%","86%", "76%","81",  "85","86%"),
    ("Clúster 8", "47","47","62", "58%", "43%","82%", "64%","50",  "82","82%"),
    ("Clúster 9", "69","47","79", "43%", "79%","27%", "28%","59",  "100","27%"),
    ("Clúster 10","71","23","63", "47%", "100%","32%","25%","30",  "100","32%"),
]

ROW_H4   = Inches(0.46)
TBL_T4   = Inches(2.2)
TBL_L4   = Inches(0.1)

# Define which cols get dark bar (high %) vs normal
dark_cols = [3, 4, 5, 6, 9]  # 0-indexed after name

# Value columns positions (matching sub_cols above)
val_lefts = [lft for _, lft in sub_cols]

for r_i, row in enumerate(clusters):
    ry = TBL_T4 + r_i * ROW_H4
    bg = RGBColor(0x0E,0x52,0xAA) if r_i % 2 == 0 else RGBColor(0x0A,0x44,0x94)
    add_rect(s4, TBL_L4, ry, SLIDE_W - TBL_L4, ROW_H4, fill=bg)

    # Row number
    txb(s4, str(r_i+1),
        TBL_L4 + Inches(.05), ry + Inches(.1),
        Inches(.25), ROW_H4,
        size=Pt(9), bold=True, color=WHITE)

    # Cluster name
    txb(s4, row[0],
        TBL_L4 + Inches(.35), ry + Inches(.1),
        Inches(1.55), ROW_H4,
        size=Pt(9), bold=True, color=WHITE)

    # Values: cols 1-3 = numeric plain, cols 4-10 = bar-style or plain
    plain_vals = [(row[1], Inches(2.05)), (row[2], Inches(2.77)), (row[3], Inches(3.5))]
    for val, lft in plain_vals:
        txb(s4, val, lft, ry + Inches(.1), Inches(.65), ROW_H4,
            size=Pt(9), color=WHITE, align=PP_ALIGN.CENTER)

    # Bar-style columns (dark bar behind %)
    bar_vals = [
        (row[4],  Inches(4.45), Inches(1.0)),
        (row[5],  Inches(5.62), Inches(1.0)),
        (row[6],  Inches(6.72), Inches(1.0)),
        (row[7],  Inches(8.08), Inches(.9)),
        (row[8],  Inches(9.08), Inches(.7)),
        (row[9],  Inches(10.15),Inches(.85)),
        (row[10], Inches(11.25),Inches(.9)),
    ]
    for val, lft, bw in bar_vals:
        # Dark bar
        add_rect(s4, lft, ry + Inches(.09), bw, ROW_H4 - Inches(.18),
                 fill=RGBColor(0x04, 0x1A, 0x55))
        txb(s4, val, lft, ry + Inches(.09), bw, ROW_H4 - Inches(.1),
            size=Pt(8.5), color=WHITE, align=PP_ALIGN.CENTER)

txb(s4, "Banco Caja Social  |  9",
    Inches(10.5), Inches(7.1), Inches(2.7), Inches(0.3),
    size=Pt(8), color=WHITE, align=PP_ALIGN.RIGHT, bold=True)

# ── SAVE ─────────────────────────────────────────────────────────────────────
# ── SLIDE 5 ─────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(blank_layout)
add_rect(s5, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)

txb(s5, "Consideraciones metodológicas incorporadas en las MIP",
    Inches(.45), Inches(.2), Inches(12.2), Inches(.65),
    size=Pt(22), bold=True, color=NAVY)
txb(s5, "Criterios aplicados a todos los países y años procesados",
    Inches(.45), Inches(.9), Inches(9.8), Inches(.35),
    size=Pt(11), color=GRAY)

add_rect(s5, Inches(.45), Inches(1.45), Inches(12.45), Inches(.55), fill=LIGHTBLUE)
txb(s5, "Z nacional/doméstica + CI importado separado + validación macro",
    Inches(.65), Inches(1.55), Inches(11.9), Inches(.35),
    size=Pt(14), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

items = [
    ("1", "Precios básicos", "Se privilegia la lectura a precios básicos. Cuando la fuente publica usos a comprador, se documenta el ajuste o aproximación."),
    ("2", "Nacional vs. importado", "Z se construye con consumo intermedio nacional/doméstico. El CI importado queda fuera de Z y en hoja separada."),
    ("3", "Consistencias macro", "Se chequea oferta = demanda y VA: g = CI nacional + CI importado + valor agregado."),
    ("4", "Actividades a reconsiderar", "La demanda final residual negativa queda marcada como revisión de clasificación, precios o apertura de uso."),
    ("5", "Vector de trabajo", "Cuando la fuente trae ocupaciones, se calculan multiplicadores de empleo; si no existe, no se imputa artificialmente."),
]

lefts = [Inches(.55), Inches(4.65), Inches(8.75)]
tops = [Inches(2.35), Inches(4.55)]
box_w = Inches(3.75)
box_h = Inches(1.55)
for idx, (num, title, body) in enumerate(items):
    col = idx % 3
    row = idx // 3
    l = lefts[col]
    t = tops[row]
    if idx == 4:
        l = Inches(4.65)
    add_rounded_rect(s5, l, t, box_w, box_h, fill=WHITE,
                     line_color=RGBColor(0xB0, 0xCC, 0xEE), line_width=Pt(1.0))
    add_rect(s5, l + Inches(.18), t + Inches(.18), Inches(.38), Inches(.38), fill=NAVY)
    txb(s5, num, l + Inches(.18), t + Inches(.13), Inches(.38), Inches(.45),
        size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s5, title, l + Inches(.68), t + Inches(.16), box_w - Inches(.85), Inches(.32),
        size=Pt(11), bold=True, color=NAVY)
    txb(s5, body, l + Inches(.2), t + Inches(.62), box_w - Inches(.4), Inches(.78),
        size=Pt(8.5), color=BLACK)

txb(s5, "Salida en Excel: hojas Z_flujos, ci_importado, demanda_final y diagnosticos_macro por país/año.",
    Inches(.55), Inches(6.65), Inches(10.5), Inches(.35),
    size=Pt(8.5), color=GRAY, italic=True)
txb(s5, "Banco Caja Social  |  Metodología MIP",
    Inches(9.3), Inches(7.1), Inches(3.6), Inches(0.3),
    size=Pt(8), color=NAVY, align=PP_ALIGN.RIGHT, bold=True)

out = "output/entregables/BancoCajaSocial_Presentacion.pptx"
prs.save(out)
print(f"[OK] {out}")
