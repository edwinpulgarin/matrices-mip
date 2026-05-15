# -*- coding: utf-8 -*-
"""Genera la PPT metodologica del proyecto CEPAL/MIP."""

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "output" / "entregables"
REPO = ROOT / "output" / "repositorio_matrices_mip"
VALIDACION = ROOT / "output" / "tablas" / "validacion_matematica_mip.xlsx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x00, 0x3F, 0x7F)
TEAL = RGBColor(0x00, 0x7B, 0xBD)
GREEN = RGBColor(0x2D, 0xBC, 0x8A)
AMBER = RGBColor(0xE8, 0xA0, 0x20)
INK = RGBColor(0x17, 0x24, 0x3A)
MUTED = RGBColor(0x5B, 0x66, 0x78)
LIGHT = RGBColor(0xF5, 0xF8, 0xFB)
PALE_BLUE = RGBColor(0xE8, 0xF4, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD9, 0xE3, 0xEF)


def add_shape(slide, shape_type, left, top, width, height, fill=None, line=None, line_width=Pt(0.75)):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_width
    return shape


def rect(slide, left, top, width, height, fill=None, line=None, line_width=Pt(0.75)):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill, line, line_width)


def round_rect(slide, left, top, width, height, fill=None, line=None, line_width=Pt(0.75)):
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill, line, line_width)
    shape.adjustments[0] = 0.08
    return shape


def tx(slide, text, left, top, width, height, size=12, bold=False, color=INK,
       align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def footer(slide, page):
    tx(slide, "CEPAL · Matrices insumo-producto · Pipeline MIP V2",
       Inches(0.55), Inches(7.12), Inches(7.8), Inches(0.22), 7.5, color=MUTED)
    tx(slide, f"{page}",
       Inches(12.3), Inches(7.12), Inches(0.45), Inches(0.22), 7.5, bold=True,
       color=NAVY, align=PP_ALIGN.RIGHT)


def title(slide, heading, subheading=None):
    tx(slide, heading, Inches(0.55), Inches(0.32), Inches(11.9), Inches(0.55),
       22, bold=True, color=NAVY)
    if subheading:
        tx(slide, subheading, Inches(0.55), Inches(0.92), Inches(10.8), Inches(0.35),
           10.5, color=MUTED)


def metric(slide, value, label, left, top, accent=TEAL):
    tx(slide, str(value), left, top, Inches(1.7), Inches(0.5), 28, bold=True, color=accent,
       align=PP_ALIGN.CENTER)
    tx(slide, label, left - Inches(0.1), top + Inches(0.55), Inches(1.9), Inches(0.42),
       8.6, color=MUTED, align=PP_ALIGN.CENTER)


def load_data():
    index = pd.read_csv(REPO / "indice_matrices.csv")
    validation = pd.read_excel(VALIDACION)
    summary = (
        index.groupby("pais")["anio"]
        .agg(["count", "min", "max"])
        .reset_index()
        .sort_values("pais")
    )
    return index, validation, summary


def add_small_bar(slide, left, top, width, color):
    rect(slide, left, top, width, Inches(0.06), fill=color)


def add_principle(slide, num, head, body, left, top, color):
    round_rect(slide, left, top, Inches(3.78), Inches(1.32), fill=WHITE, line=LINE)
    rect(slide, left, top, Inches(0.11), Inches(1.32), fill=color)
    tx(slide, num, left + Inches(0.23), top + Inches(0.18), Inches(0.35), Inches(0.28),
       12, bold=True, color=color, align=PP_ALIGN.CENTER)
    tx(slide, head, left + Inches(0.72), top + Inches(0.18), Inches(2.85), Inches(0.28),
       10.5, bold=True, color=NAVY)
    tx(slide, body, left + Inches(0.72), top + Inches(0.56), Inches(2.85), Inches(0.56),
       8.2, color=INK)


def add_flow_node(slide, label, body, left, top, fill, width=Inches(2.15)):
    round_rect(slide, left, top, width, Inches(1.05), fill=fill, line=None)
    tx(slide, label, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.24),
       10.2, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(slide, body, left + Inches(0.18), top + Inches(0.48), width - Inches(0.36), Inches(0.36),
       7.6, color=WHITE, align=PP_ALIGN.CENTER)


def arrow(slide, left, top, width=Inches(0.55)):
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.28), fill=AMBER, line=None)


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    index, validation, country_summary = load_data()
    total = len(index)
    structural_ok = int((validation["validacion_estructural"] == "OK").sum())
    names_ok = int(validation["nombres_sector_ok"].sum())

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT)
    rect(slide, 0, 0, Inches(0.28), SLIDE_H, fill=NAVY)
    tx(slide, "CEPAL · MIP AMERICA LATINA", Inches(0.62), Inches(0.55), Inches(5), Inches(0.28),
       9.5, bold=True, color=TEAL)
    tx(slide, "Repositorio de matrices\ninsumo-producto", Inches(0.62), Inches(1.35),
       Inches(7.25), Inches(1.7), 34, bold=True, color=NAVY)
    tx(slide, "Construccion comparable para Argentina, Brasil, Mexico y Uruguay desde COU/MIP oficiales y bases CEPAL.",
       Inches(0.66), Inches(3.18), Inches(6.4), Inches(0.78), 15, color=INK)
    add_small_bar(slide, Inches(0.66), Inches(4.18), Inches(2.6), AMBER)
    metric(slide, total, "matrices pais-anio", Inches(0.7), Inches(4.78), TEAL)
    metric(slide, 4, "paises cubiertos", Inches(2.7), Inches(4.78), GREEN)
    metric(slide, "1997-2021", "rango temporal", Inches(4.7), Inches(4.78), AMBER)
    round_rect(slide, Inches(8.35), Inches(1.1), Inches(4.25), Inches(4.95), fill=WHITE, line=LINE)
    tx(slide, "Entregables principales", Inches(8.72), Inches(1.45), Inches(3.5), Inches(0.35),
       14, bold=True, color=NAVY)
    bullets = [
        "Excel por pais y anio con Z, A, L, Ghosh, VA, produccion y balances.",
        "Consumo intermedio importado separado de Z.",
        "Validacion matematica y diagnostica reproducible.",
        "Repositorio local versionado y paquete para Google Drive.",
    ]
    for i, item in enumerate(bullets):
        y = Inches(2.05 + 0.72 * i)
        add_shape(slide, MSO_SHAPE.OVAL, Inches(8.78), y + Inches(0.03), Inches(0.12), Inches(0.12),
                  fill=GREEN if i % 2 else TEAL, line=None)
        tx(slide, item, Inches(9.05), y - Inches(0.02), Inches(3.1), Inches(0.45), 9.2, color=INK)
    tx(slide, "Mayo 2026", Inches(8.72), Inches(5.52), Inches(3.4), Inches(0.28),
       8.8, color=MUTED)
    footer(slide, 1)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    title(slide, "Alcance del repositorio", "Cobertura pais-anio y fuentes usadas en la construccion.")
    headers = ["Pais", "Matrices", "Anios", "Fuente principal"]
    sources = {
        "Argentina": "INDEC + CEPAL + MIPAr97",
        "Brasil": "IBGE + CEPAL",
        "Mexico": "INEGI + CEPAL",
        "Uruguay": "BCU + CEPAL",
    }
    left, top = Inches(0.7), Inches(1.6)
    widths = [Inches(2.2), Inches(1.45), Inches(2.45), Inches(4.5)]
    rect(slide, left, top, sum(widths), Inches(0.42), fill=NAVY)
    x = left
    for h, w in zip(headers, widths):
        tx(slide, h, x + Inches(0.1), top + Inches(0.11), w - Inches(0.2), Inches(0.2),
           8.6, bold=True, color=WHITE)
        x += w
    for r, row in enumerate(country_summary.itertuples(index=False), start=1):
        y = top + Inches(0.42 * r)
        bg = LIGHT if r % 2 else WHITE
        rect(slide, left, y, sum(widths), Inches(0.42), fill=bg, line=LINE, line_width=Pt(0.35))
        vals = [row.pais, row.count, f"{row.min}-{row.max}", sources.get(row.pais, "Fuente oficial")]
        x = left
        for val, w in zip(vals, widths):
            tx(slide, str(val), x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), Inches(0.2),
               8.3, bold=(w == widths[0]), color=INK)
            x += w
    round_rect(slide, Inches(0.7), Inches(4.45), Inches(5.3), Inches(1.3), fill=PALE_BLUE, line=None)
    tx(slide, "Estructura navegable", Inches(1.0), Inches(4.75), Inches(4.7), Inches(0.28),
       12, bold=True, color=NAVY)
    tx(slide, "repositorio_matrices_mip/{Pais}/{Anio}/MIP_{Pais}_{Anio}.xlsx\nindice_matrices.xlsx + METODOLOGIA.md + validacion_matematica_mip.xlsx",
       Inches(1.0), Inches(5.13), Inches(4.7), Inches(0.48), 8.5, color=INK)
    round_rect(slide, Inches(6.45), Inches(4.45), Inches(5.95), Inches(1.3), fill=WHITE, line=LINE)
    tx(slide, "Regla de publicacion", Inches(6.8), Inches(4.75), Inches(5.2), Inches(0.28),
       12, bold=True, color=NAVY)
    tx(slide, "Toda matriz publicada debe tener nombres de sectores economicos en filas y columnas; los codigos pueden quedar como prefijo, no como reemplazo.",
       Inches(6.8), Inches(5.13), Inches(5.2), Inches(0.48), 8.7, color=INK)
    footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT)
    title(slide, "Decisiones metodologicas transversales", "Criterios aplicados a todos los paises y anios procesados.")
    principles = [
        ("1", "Precios basicos", "Se trabaja a precios basicos cuando la fuente lo permite; cualquier aproximacion queda documentada.", TEAL),
        ("2", "Z nacional/domestica", "La matriz Z contiene solo consumo intermedio nacional o domestico.", GREEN),
        ("3", "CI importado separado", "El consumo intermedio importado queda fuera de Z y se publica como vector propio.", AMBER),
        ("4", "Cierres macro", "Se revisa oferta = demanda y g = CI nacional + CI importado + valor agregado.", NAVY),
        ("5", "Actividades a reconsiderar", "Los residuales negativos se conservan como diagnostico sectorial, no se ocultan.", TEAL),
        ("6", "Vector de trabajo", "Los multiplicadores de empleo se calculan solo cuando la fuente trae ocupaciones.", GREEN),
    ]
    xs = [Inches(0.68), Inches(4.78), Inches(8.88)]
    ys = [Inches(1.72), Inches(3.42)]
    for i, item in enumerate(principles):
        num, head, body, color = item
        add_principle(slide, num, head, body, xs[i % 3], ys[i // 3], color)
    round_rect(slide, Inches(0.68), Inches(5.55), Inches(12.0), Inches(0.76), fill=WHITE, line=LINE)
    tx(slide, "Salida en Excel", Inches(0.95), Inches(5.77), Inches(1.3), Inches(0.22),
       9.5, bold=True, color=NAVY)
    tx(slide, "Z_MIP, A_coef_tecnicos, L_leontief, B/G Ghosh, g_produccion, W_valor_agregado, f_demanda_final, CI_importado, multiplicadores y balances_sectoriales.",
       Inches(2.25), Inches(5.77), Inches(10.0), Inches(0.24), 8.3, color=INK)
    footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    title(slide, "Conversion COU a matriz insumo-producto", "Supuesto de tecnologia de industria para convertir tablas rectangulares en matrices sector-sector.")
    y = Inches(2.08)
    nodes = [
        ("V", "produccion por industria y producto", NAVY),
        ("U nacional", "uso intermedio domestico", TEAL),
        ("D", "estructura producto-industria", GREEN),
        ("Z", "flujos intermedios sector-sector", NAVY),
        ("A, L", "coeficientes e inversa de Leontief", AMBER),
    ]
    x = Inches(0.72)
    for i, (label, body, color) in enumerate(nodes):
        add_flow_node(slide, label, body, x, y, color)
        if i < len(nodes) - 1:
            arrow(slide, x + Inches(2.28), y + Inches(0.38))
        x += Inches(2.72)
    round_rect(slide, Inches(0.9), Inches(4.25), Inches(5.65), Inches(1.18), fill=PALE_BLUE, line=None)
    tx(slide, "Identidades de transformacion", Inches(1.2), Inches(4.55), Inches(4.9), Inches(0.25),
       11.5, bold=True, color=NAVY)
    tx(slide, "D = V · diag(q)^-1\nZ = D · U_nacional\nA = Z · diag(g)^-1     L = (I - A)^-1",
       Inches(1.2), Inches(4.92), Inches(4.9), Inches(0.45), 9.1, color=INK)
    round_rect(slide, Inches(6.95), Inches(4.25), Inches(5.45), Inches(1.18), fill=LIGHT, line=LINE)
    tx(slide, "Lectura economica", Inches(7.25), Inches(4.55), Inches(4.8), Inches(0.25),
       11.5, bold=True, color=NAVY)
    tx(slide, "La MIP resultante representa compras y ventas intermedias entre sectores economicos, dejando importaciones intermedias fuera de Z para no inflar encadenamientos domesticos.",
       Inches(7.25), Inches(4.92), Inches(4.8), Inches(0.45), 8.8, color=INK)
    footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT)
    title(slide, "Separacion nacional/importado", "La prioridad es que Z mida encadenamientos domesticos; el componente importado queda visible y separado.")
    rect(slide, Inches(0.9), Inches(2.05), Inches(2.4), Inches(0.72), fill=NAVY)
    tx(slide, "U total", Inches(0.9), Inches(2.25), Inches(2.4), Inches(0.25), 13,
       bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    arrow(slide, Inches(3.55), Inches(2.27), Inches(0.65))
    rect(slide, Inches(4.45), Inches(1.55), Inches(3.0), Inches(0.72), fill=GREEN)
    tx(slide, "U nacional", Inches(4.45), Inches(1.75), Inches(3.0), Inches(0.25),
       13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide, Inches(4.45), Inches(2.72), Inches(3.0), Inches(0.72), fill=AMBER)
    tx(slide, "U importada / CI importado", Inches(4.45), Inches(2.92), Inches(3.0), Inches(0.25),
       12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    arrow(slide, Inches(7.72), Inches(1.77), Inches(0.65))
    rect(slide, Inches(8.62), Inches(1.55), Inches(3.0), Inches(0.72), fill=NAVY)
    tx(slide, "Z_MIP", Inches(8.62), Inches(1.75), Inches(3.0), Inches(0.25),
       13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    arrow(slide, Inches(7.72), Inches(2.94), Inches(0.65))
    rect(slide, Inches(8.62), Inches(2.72), Inches(3.0), Inches(0.72), fill=WHITE, line=AMBER)
    tx(slide, "CI_importado", Inches(8.62), Inches(2.92), Inches(3.0), Inches(0.25),
       12, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    round_rect(slide, Inches(0.9), Inches(4.45), Inches(5.55), Inches(1.05), fill=WHITE, line=LINE)
    tx(slide, "Si solo existe vector de importaciones por producto", Inches(1.2), Inches(4.72),
       Inches(5.0), Inches(0.22), 10.5, bold=True, color=NAVY)
    tx(slide, "Se asigna proporcionalmente: participacion_importada_p = M_p / (produccion_domestica_p + M_p).",
       Inches(1.2), Inches(5.08), Inches(5.0), Inches(0.28), 8.5, color=INK)
    round_rect(slide, Inches(6.95), Inches(4.45), Inches(5.45), Inches(1.05), fill=WHITE, line=LINE)
    tx(slide, "Cierre de valor agregado", Inches(7.25), Inches(4.72),
       Inches(4.9), Inches(0.22), 10.5, bold=True, color=NAVY)
    tx(slide, "g_j = compras_intermedias_nacionales_j + CI_importado_j + valor_agregado_j",
       Inches(7.25), Inches(5.08), Inches(4.9), Inches(0.28), 8.5, color=INK)
    footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    title(slide, "Validacion y alertas economicas", "Las pruebas estructurales bloquean errores; las alertas diagnosticas orientan revision economica.")
    metric(slide, f"{structural_ok}/{total}", "validacion estructural OK", Inches(0.95), Inches(1.78), TEAL)
    metric(slide, f"{names_ok}/{total}", "nombres sectoriales OK", Inches(3.15), Inches(1.78), GREEN)
    metric(slide, total, "archivos en repositorio", Inches(5.35), Inches(1.78), AMBER)
    round_rect(slide, Inches(0.8), Inches(3.35), Inches(5.9), Inches(1.6), fill=PALE_BLUE, line=None)
    tx(slide, "Chequeos estructurales", Inches(1.1), Inches(3.7), Inches(5.3), Inches(0.3),
       12, bold=True, color=NAVY)
    tx(slide, "- Z, A y L cuadradas y alineadas\n- A = Z / g\n- (I - A)L = I\n- Nombres sectoriales en filas y columnas",
       Inches(1.1), Inches(4.08), Inches(5.3), Inches(0.72), 8.8, color=INK)
    round_rect(slide, Inches(7.05), Inches(1.55), Inches(5.45), Inches(3.4), fill=LIGHT, line=LINE)
    tx(slide, "Chequeos diagnosticos macro", Inches(7.4), Inches(1.92), Inches(4.8), Inches(0.3),
       12, bold=True, color=NAVY)
    tx(slide, "Oferta = demanda:\ng_i = ventas_intermedias_nacionales_i + demanda_final_i",
       Inches(7.4), Inches(2.35), Inches(4.8), Inches(0.58), 8.8, color=INK)
    tx(slide, "Valor agregado:\ng_j = compras_intermedias_nacionales_j + CI_importado_j + W_j",
       Inches(7.4), Inches(3.15), Inches(4.8), Inches(0.58), 8.8, color=INK)
    tx(slide, "Actividades a reconsiderar: sectores con demanda final residual negativa o VA residual negativo quedan marcados para revision.",
       Inches(7.4), Inches(4.02), Inches(4.8), Inches(0.54), 8.6, color=INK)
    footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT)
    title(slide, "Repositorio compartible del proyecto", "Distribucion por Google Drive y versionamiento local del paquete de matrices.")
    steps = [
        ("1", "Empaquetar", "ZIP limpio del repositorio, sin depender de rutas locales."),
        ("2", "Subir a Drive", "Carpeta compartida: Repositorio MIP CEPAL."),
        ("3", "Acceso de lectura", "Lector para consulta; Comentador para revision metodologica."),
        ("4", "Actualizar", "Regenerar ZIP despues de cada corrida validada."),
    ]
    for i, (num, head, body) in enumerate(steps):
        left = Inches(0.8 + i * 3.05)
        round_rect(slide, left, Inches(1.75), Inches(2.55), Inches(1.72), fill=WHITE, line=LINE)
        rect(slide, left + Inches(0.22), Inches(2.0), Inches(0.42), Inches(0.42), fill=NAVY)
        tx(slide, num, left + Inches(0.22), Inches(2.08), Inches(0.42), Inches(0.18),
           10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tx(slide, head, left + Inches(0.78), Inches(2.0), Inches(1.5), Inches(0.28),
           10.5, bold=True, color=NAVY)
        tx(slide, body, left + Inches(0.28), Inches(2.52), Inches(2.0), Inches(0.55),
           8.1, color=INK)
    round_rect(slide, Inches(1.05), Inches(4.45), Inches(5.4), Inches(1.18), fill=WHITE, line=LINE)
    tx(slide, "Paquete para el equipo", Inches(1.35), Inches(4.76), Inches(4.8), Inches(0.26),
       11.5, bold=True, color=NAVY)
    tx(slide, "Repositorio_MIP_CEPAL_equipo.zip\nCEPAL_MIP_Metodologia.pptx\nvalidacion_matematica_mip.xlsx",
       Inches(1.35), Inches(5.1), Inches(4.8), Inches(0.46), 8.8, color=INK)
    round_rect(slide, Inches(6.95), Inches(4.45), Inches(5.2), Inches(1.18), fill=PALE_BLUE, line=None)
    tx(slide, "Si quieren Git mas adelante", Inches(7.25), Inches(4.76), Inches(4.6), Inches(0.26),
       11.5, bold=True, color=NAVY)
    tx(slide, "El repositorio local queda versionado. Para GitHub, usar acceso de lectura o colaboradores segun corresponda.",
       Inches(7.25), Inches(5.12), Inches(4.6), Inches(0.4), 8.8, color=INK)
    footer(slide, 7)

    out = OUT / "CEPAL_MIP_Metodologia.pptx"
    prs.save(out)
    print(f"[OK] {out}")


if __name__ == "__main__":
    main()
